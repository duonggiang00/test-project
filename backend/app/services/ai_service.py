import time
from app.db.session import SessionLocal
from app.models.material import StudyMaterial
from app.models.exam import Question, Option
from app.models.topic import Topic
from app.models.flashcard import FlashcardDeck, Flashcard
from app.models.document_chunk import DocumentChunk
import os

def mock_process_document_and_generate_questions(material_id: str):
    # Simulate processing time
    time.sleep(5)
    
    with SessionLocal() as db:
        material = db.query(StudyMaterial).filter(StudyMaterial.id == material_id).first()
        if not material:
            return
            
        # 1. Process Document and Create Chunks
        try:
            content = ""
            if material.file_path.lower().endswith(".pdf"):
                import pdfplumber
                with pdfplumber.open(material.file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            content += page_text + "\n\n"
            elif material.file_path.lower().endswith(".pptx"):
                import pptx
                prs = pptx.Presentation(material.file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            content += shape.text + "\n\n"
            elif material.file_path.lower().endswith(".docx"):
                import docx
                doc = docx.Document(material.file_path)
                for para in doc.paragraphs:
                    if para.text:
                        content += para.text + "\n\n"
            else:
                with open(material.file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            
            
            # Simple chunking by paragraph (split by double newline)
            paragraphs = [p.strip() for p in content.split("\n\n") if len(p.strip()) > 10]
            
            if not paragraphs:
                paragraphs = [content[:1000]] # Fallback if no newlines
                
            for para in paragraphs:
                chunk = DocumentChunk(
                    material_id=material.id,
                    content=para[:2000], # Limit size
                    embedding=[0.1] * 1536 # Mock embedding
                )
                db.add(chunk)
            db.commit()
        except Exception as e:
            print(f"Failed to chunk document: {e}")
            material.ai_status = 'failed'
            db.commit()
            return
            
        # Mock generating questions
        q1 = Question(
            material_id=material.id,
            content=f"What is the main topic of {material.title}?",
            is_ai_generated=True,
            points=1
        )
        db.add(q1)
        db.flush() # get q1.id
        
        db.add(Option(question_id=q1.id, content="AI", is_correct=True))
        db.add(Option(question_id=q1.id, content="Blockchain", is_correct=False))
        
        q2 = Question(
            material_id=material.id,
            content="Which statement is true based on the document?",
            is_ai_generated=True,
            points=1
        )
        db.add(q2)
        db.flush()
        
        db.add(Option(question_id=q2.id, content="This is true", is_correct=True))
        db.add(Option(question_id=q2.id, content="This is false", is_correct=False))
        
        material.ai_status = "completed"
        db.commit()

def mock_generate_topic_kit(material_id: str, topic_id: str):
    """
    Mock function to simulate AI generating a Topic Brief and Flashcards from a Study Material.
    """
    time.sleep(3) # Simulate AI processing time
    
    with SessionLocal() as db:
        material = db.query(StudyMaterial).filter(StudyMaterial.id == material_id).first()
        topic = db.query(Topic).filter(Topic.id == topic_id).first()
        
        if not material or not topic:
            return

        # 1. Generate Topic Brief
        topic.brief_content = f"# {topic.name}\n\nĐây là bài tóm tắt kiến thức được tạo tự động bởi AI từ tài liệu: **{material.title}**.\n\n## 1. Khái niệm cơ bản\nNội dung khái niệm...\n\n## 2. Các điểm trọng tâm\n- Điểm 1\n- Điểm 2\n"
        topic.brief_ai_generated = True

        # 2. Generate Flashcard Deck
        deck = FlashcardDeck(
            topic_id=topic.id,
            title=f"Flashcards: {material.title}",
            description="Bộ thẻ ghi nhớ tự động tạo từ tài liệu."
        )
        db.add(deck)
        db.flush()

        # 3. Generate Flashcards
        cards_data = [
            ("AI là gì?", "Trí tuệ nhân tạo (AI) là khả năng của máy tính bắt chước các chức năng nhận thức của con người."),
            ("Spaced Repetition là gì?", "Kỹ thuật ôn tập ngắt quãng giúp ghi nhớ dài hạn bằng cách tăng dần thời gian giữa các lần ôn tập."),
            ("Mô hình Dữ liệu (Data Model) là gì?", "Cách cấu trúc và tổ chức dữ liệu trong cơ sở dữ liệu.")
        ]

        for i, (front, back) in enumerate(cards_data):
            card = Flashcard(
                deck_id=deck.id,
                front_content=front,
                back_content=back,
                order_index=i
            )
            db.add(card)

        db.commit()
