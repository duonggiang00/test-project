from __future__ import annotations

import re
from pathlib import Path


MOCK_EMBEDDING_DIMENSIONS = 1536
MOCK_EMBEDDING = [0.1] * MOCK_EMBEDDING_DIMENSIONS


def extract_material_text(path: Path) -> str:
    """Extract text from a supported material using the product's local toolchain."""
    suffix = path.suffix.casefold()
    if suffix == ".pdf":
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            return "\n\n".join(
                page_text
                for page in pdf.pages
                if (page_text := page.extract_text())
            )
    if suffix == ".pptx":
        import pptx

        presentation = pptx.Presentation(str(path))
        return "\n\n".join(
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
    if suffix == ".docx":
        import docx

        document = docx.Document(path)
        return "\n\n".join(
            paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()
        )
    return path.read_text(encoding="utf-8")


def chunk_material_text(content: str, max_characters: int = 2000) -> list[str]:
    """Create deterministic, whitespace-normalized chunks without provider calls."""
    normalized = re.sub(r"\s+", " ", content).strip()
    if not normalized:
        return []

    chunks: list[str] = []
    remaining = normalized
    while len(remaining) > max_characters:
        boundary = remaining.rfind(". ", 0, max_characters)
        if boundary < max_characters // 2:
            boundary = remaining.rfind(" ", 0, max_characters)
        if boundary <= 0:
            boundary = max_characters
        else:
            boundary += 1
        chunks.append(remaining[:boundary].strip())
        remaining = remaining[boundary:].strip()
    if remaining:
        chunks.append(remaining)
    return chunks


def extract_and_chunk_material(path: Path) -> tuple[str, list[str]]:
    content = extract_material_text(path)
    return content, chunk_material_text(content)
